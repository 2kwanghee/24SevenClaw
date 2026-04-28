"""가이드 PPTX 생성기 — python-pptx 기반 설치/실행 가이드 슬라이드 생성."""

import io

from pptx import Presentation  # type: ignore[import-untyped]
from pptx.dml.color import RGBColor  # type: ignore[import-untyped]
from pptx.enum.text import PP_ALIGN  # type: ignore[import-untyped]
from pptx.util import Inches, Pt  # type: ignore[import-untyped]

_BG = RGBColor(0x10, 0x18, 0x1A)
_ACCENT = RGBColor(0x10, 0xB9, 0x81)
_WHITE = RGBColor(0xFF, 0xFF, 0xFF)
_GRAY = RGBColor(0x94, 0xA3, 0xB8)
_AMBER = RGBColor(0xFB, 0xBF, 0x24)
_SKY = RGBColor(0x38, 0xBD, 0xF8)


def _set_bg(slide: object) -> None:
    fill = slide.background.fill  # type: ignore[attr-defined]
    fill.solid()
    fill.fore_color.rgb = _BG


def _add_text(
    slide: object,
    text: str,
    left: float,
    top: float,
    width: float,
    height: float,
    *,
    bold: bool = False,
    size: int = 16,
    color: RGBColor | None = None,
    align: object = PP_ALIGN.LEFT,
) -> None:
    tx_box = slide.shapes.add_textbox(left, top, width, height)  # type: ignore[attr-defined]
    tf = tx_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.bold = bold
    run.font.size = Pt(size)
    run.font.color.rgb = color if color is not None else _WHITE


def _slide_header(slide: object, title: str, slide_w: float) -> None:
    _add_text(slide, title, Inches(0.8), Inches(0.5), slide_w, Inches(0.8),
              bold=True, size=22, color=_ACCENT)


def _numbered_items(
    slide: object,
    items: list[tuple[str, str | None]],
    start_y: float = 1.5,
    body_w: float = Inches(8.5),
) -> None:
    """(label, command_or_None) 리스트를 번호 목록으로 추가."""
    y = start_y
    for label, cmd in items:
        _add_text(slide, label, Inches(0.8), Inches(y), body_w, Inches(0.55),
                  size=14, color=_WHITE)
        y += 0.55
        if cmd:
            _add_text(slide, cmd, Inches(1.2), Inches(y), Inches(7.0), Inches(0.5),
                      bold=True, size=13, color=_ACCENT)
            y += 0.6
        y += 0.1


def build_setup_guide_pptx(
    project_name: str,
    pm_slug: str,
    has_linear: bool,
    platform: str,
) -> bytes:
    """설치/실행 가이드 PPTX를 생성하여 bytes로 반환."""
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]

    slide_w = Inches(10)
    body_w = Inches(8.5)

    # ── 슬라이드 1: 표지 ──────────────────────────────────────────────────────
    s = prs.slides.add_slide(blank)
    _set_bg(s)
    _add_text(s, "ClickEye", Inches(0.8), Inches(1.0), slide_w, Inches(0.7),
              bold=True, size=13, color=_ACCENT)
    _add_text(s, project_name, Inches(0.8), Inches(1.9), slide_w, Inches(1.4),
              bold=True, size=34, color=_WHITE)
    _add_text(s, "설치 및 실행 가이드", Inches(0.8), Inches(3.3), slide_w, Inches(0.7),
              size=20, color=_GRAY)
    _add_text(s, f"플랫폼: {platform}", Inches(0.8), Inches(4.2), slide_w, Inches(0.6),
              size=14, color=_GRAY)
    if has_linear:
        _add_text(s, "Linear 연동 포함", Inches(0.8), Inches(4.9), slide_w, Inches(0.6),
                  size=13, color=_SKY)

    # ── 슬라이드 2: 다운로드 & 압축 해제 ─────────────────────────────────────
    s = prs.slides.add_slide(blank)
    _set_bg(s)
    _slide_header(s, "Step 1 — 다운로드 & 압축 해제", slide_w)
    unzip_cmd = f"unzip {project_name}.zip -d ~/projects/{project_name}"
    _numbered_items(s, [
        (f"① 프로젝트 대시보드에서  {project_name}.zip  다운로드", None),
        ("② WSL2 Ubuntu 터미널에서 압축 해제", unzip_cmd),
        ("③ 프로젝트 폴더로 이동", f"cd ~/projects/{project_name}"),
    ])
    _add_text(s, "💡  Windows 사용자: WSL2 Ubuntu 터미널에서 실행하세요 (탐색기 더블클릭 X)",
              Inches(0.8), Inches(5.8), body_w, Inches(0.6), size=12, color=_AMBER)

    # ── 슬라이드 3: .env 설정 ─────────────────────────────────────────────────
    s = prs.slides.add_slide(blank)
    _set_bg(s)
    _slide_header(s, "Step 2 — .env API 키 설정", slide_w)
    _add_text(s, ".env.example → .env 로 복사 후 아래 키를 입력하세요",
              Inches(0.8), Inches(1.35), body_w, Inches(0.55), size=14, color=_GRAY)
    _add_text(s, "cp .env.example .env", Inches(0.8), Inches(1.95), Inches(6), Inches(0.5),
              bold=True, size=13, color=_ACCENT)
    env_keys: list[str] = ["ANTHROPIC_API_KEY=sk-ant-..."]
    if has_linear:
        env_keys += [
            "LINEAR_API_KEY=lin_api_...",
            "LINEAR_TEAM_ID=<팀 UUID>",
            "WEBHOOK_SECRET=<openssl rand -hex 32>",
        ]
    y = 2.6
    for kv in env_keys:
        _add_text(s, kv, Inches(0.8), Inches(y), Inches(7), Inches(0.55),
                  bold=True, size=13, color=_ACCENT)
        y += 0.65
    _add_text(s, "📁  docs/api-keys/ 폴더에서 각 키 발급 방법을 확인하세요.",
              Inches(0.8), Inches(y + 0.2), body_w, Inches(0.6), size=12, color=_GRAY)

    if has_linear:
        # ── 슬라이드 4 (Linear): 런처 실행 ──────────────────────────────────
        s = prs.slides.add_slide(blank)
        _set_bg(s)
        _slide_header(s, "Step 3 — 런처 실행", slide_w)
        _add_text(s, "start.sh 하나로 터널·Webhook·Watcher가 자동 기동됩니다.",
                  Inches(0.8), Inches(1.35), body_w, Inches(0.55), size=14, color=_GRAY)
        _numbered_items(s, [
            ("① 터미널에서 런처 실행", "bash start.sh"),
            ("② Cloudflare Tunnel 자동 기동 → Linear Webhook 자동 등록", None),
            ("③ Webhook 서버 + Watcher 백그라운드 기동", None),
            ("④ 완료 후 스크립트 자동 종료 — 터미널을 닫아도 됩니다", None),
        ], start_y=2.0)
        _add_text(s, "💡  start.sh 완료 후 터미널을 닫아도 파이프라인은 계속 실행됩니다.",
                  Inches(0.8), Inches(5.8), body_w, Inches(0.6), size=12, color=_ACCENT)

        # ── 슬라이드 5 (Linear): 서비스 영구 등록 ────────────────────────────
        s = prs.slides.add_slide(blank)
        _set_bg(s)
        _slide_header(s, "Step 4 — 서비스 영구 등록 (권장)", slide_w)
        _add_text(s, "WSL2 재시작 후에도 서비스가 자동 복구되도록 systemd에 등록합니다.",
                  Inches(0.8), Inches(1.35), body_w, Inches(0.55), size=14, color=_GRAY)
        _numbered_items(s, [
            ("① systemd 사용자 서비스 등록", "bash scripts/install-service.sh"),
            ("② 터미널 종료·WSL2 재시작 후에도 자동 유지 확인", None),
            ("③ (선택) Windows 재부팅 후 WSL2 자동 기동 — 관리자 PowerShell에서:",
             "powershell -ExecutionPolicy Bypass -File scripts\\setup-autostart.ps1"),
        ], start_y=2.0)
        _add_text(
            s,
            "⚠️  Cloudflare 무료 터널은 WSL2 재시작 시 URL이 바뀝니다."
            " 재시작 후 bash start.sh를 다시 실행해 Webhook을 재등록하세요.",
            Inches(0.8), Inches(5.5), body_w, Inches(0.8), size=12, color=_AMBER,
        )

        # ── 슬라이드 6 (Linear): AI Team → Linear 자동화 ─────────────────────
        s = prs.slides.add_slide(blank)
        _set_bg(s)
        _slide_header(s, "Step 5 — AI Team → Linear 자동화", slide_w)
        _add_text(s, "셋업 완료 후 AI Team 메뉴에서 작업 요청을 등록하세요.",
                  Inches(0.8), Inches(1.35), body_w, Inches(0.55), size=14, color=_GRAY)
        _numbered_items(s, [
            ("① ClickEye 웹 → 프로젝트 → AI Team 메뉴 진입", None),
            ("② '새 작업 요청' 클릭 → 제목·설명 입력 → '생성 & 분해'", None),
            ("③ 서브태스크 확인 → '배정 확정' 클릭", None),
            ("④ 배정 완료 → Linear 이슈 자동 등록 (서브태스크별 1개)",
             "예: [backend] 인증 API, [frontend] 로그인 UI, [qa] 테스트"),
            ("⑤ ClickEye 웹에서 태스크 승인(Todo) → 로컬 Claude가 자동으로 작업 시작", None),
        ], start_y=2.0)

    else:
        # ── 슬라이드 4 (No-Linear): 런처 실행 ────────────────────────────────
        s = prs.slides.add_slide(blank)
        _set_bg(s)
        _slide_header(s, "Step 3 — 런처 실행", slide_w)
        _numbered_items(s, [
            ("① 터미널에서 런처 스크립트 실행",
             "bash start.sh"),
            ("② Node.js · Claude Code 자동 점검", None),
            ("③ .env 검증 → 누락 키 대화형 안내", None),
            ("④ 완료 후 스크립트 자동 종료 — 터미널을 닫아도 됩니다", None),
        ], start_y=1.5)
        _add_text(s, "💡  start.sh 완료 후 터미널을 닫아도 됩니다.",
                  Inches(0.8), Inches(5.8), body_w, Inches(0.6), size=12, color=_ACCENT)

        # ── 슬라이드 5 (No-Linear): 서비스 영구 등록 ─────────────────────────
        s = prs.slides.add_slide(blank)
        _set_bg(s)
        _slide_header(s, "Step 4 — 서비스 영구 등록 (권장)", slide_w)
        _add_text(s, "WSL2 재시작 후에도 서비스가 자동 복구되도록 systemd에 등록합니다.",
                  Inches(0.8), Inches(1.35), body_w, Inches(0.55), size=14, color=_GRAY)
        _numbered_items(s, [
            ("① systemd 사용자 서비스 등록", "bash scripts/install-service.sh"),
            ("② 터미널 종료·WSL2 재시작 후에도 자동 유지 확인", None),
            ("③ (선택) Windows 재부팅 후 WSL2 자동 기동 — 관리자 PowerShell에서:",
             "powershell -ExecutionPolicy Bypass -File scripts\\setup-autostart.ps1"),
        ], start_y=2.0)

    # ── 마지막 슬라이드: 참고 문서 & 다음 단계 ───────────────────────────────
    s = prs.slides.add_slide(blank)
    _set_bg(s)
    _slide_header(s, "참고 문서 & 다음 단계", slide_w)
    if has_linear:
        refs = [
            ("docs/api-keys/", "각 API 키 발급 상세 가이드"),
            ("docs/WEBHOOK_SETUP.md", "Linear Webhook / 터널 / 폴링 설정 상세"),
            ("scripts/install-service.sh", "systemd 서비스 등록 (영구 유지)"),
            ("ClickEye 웹 → AI Team", "작업 요청 등록 → Linear 이슈 자동 생성"),
        ]
    else:
        refs = [
            ("docs/api-keys/", "각 API 키 발급 상세 가이드"),
            ("scripts/install-service.sh", "systemd 서비스 등록 (영구 유지)"),
            ("프로젝트 대시보드", "설정 변경 / ZIP 재다운로드"),
        ]

    for i, (path, desc) in enumerate(refs):
        _add_text(s, f"📄  {path}",
                  Inches(0.8), Inches(1.5 + i * 1.1), Inches(3.8), Inches(0.6),
                  bold=True, size=13, color=_ACCENT)
        _add_text(s, desc,
                  Inches(4.8), Inches(1.5 + i * 1.1), Inches(4.8), Inches(0.6),
                  size=13, color=_GRAY)
    _add_text(s, "문의 / 피드백: 대시보드 내 Contact 페이지",
              Inches(0.8), Inches(6.3), body_w, Inches(0.6), size=12, color=_GRAY)

    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()
